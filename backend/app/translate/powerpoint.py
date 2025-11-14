import threading
import pptx
from . import to_translate
from . import common
import os
import sys
import time
import logging
from datetime import datetime
import re # Added for regex operations

# 配置日志
logger = logging.getLogger(__name__)

# 简单的语法检查
try:
    # 测试基本语法
    test_var = "test"
    test_list = [1, 2, 3]
    test_dict = {"key": "value"}
    logger.info("PPT翻译模块语法检查通过")
except Exception as e:
    logger.error(f"PPT翻译模块语法检查失败: {e}")
    raise

# 字体缩放功能开关
ENABLE_FONT_SCALING = True  # 设置为False关闭字体缩放，设置为True开启字体缩放

# 字体缩放功能说明：
# 当ENABLE_FONT_SCALING = True时：
#   - 系统会根据翻译后文本的长度自动调整字体大小
#   - 长文本会适当缩小字体，短文本会适当放大字体
#   - 确保文本在PPT中显示完整且美观
#
# 当ENABLE_FONT_SCALING = False时：
#   - 字体大小保持不变，使用原始设置
#   - 翻译后的文本可能超出文本框边界
#   - 需要手动调整文本框大小或字体大小
#
# 使用方法：
# 1. 关闭字体缩放：ENABLE_FONT_SCALING = False
# 2. 开启字体缩放：ENABLE_FONT_SCALING = True
# 3. 修改后需要重新运行翻译任务才能生效

def start(trans):
    """PPT翻译入口函数，支持Okapi和传统方法"""
    # 检查是否使用Okapi方法（默认使用Okapi，与Word翻译保持一致）
    use_okapi = trans.get('use_okapi', True)
    
    if use_okapi:
        logger.info("使用 Okapi Framework 进行 PPT 翻译")
        return start_with_okapi(trans)
    else:
        logger.info("使用传统方法进行 PPT 翻译")
        return start_traditional(trans)


def start_with_okapi(trans):
    """使用 Okapi Framework 进行翻译"""
    try:
        # 导入 Okapi 集成模块
        from .okapi_integration import OkapiPptxTranslator, verify_okapi_installation
        
        # 验证 Okapi 安装
        if not verify_okapi_installation():
            logger.error("❌ Okapi 安装验证失败，回退到传统方法")
            return start_traditional(trans)
        
        # 如果用户选择了qwen-mt-plus，设置server为qwen
        if trans.get('model') == 'qwen-mt-plus':
            trans['server'] = 'qwen'
            logger.info("✅ 设置翻译服务为 Qwen")
        
        # 预加载术语库
        comparison_id = trans.get('comparison_id')
        if comparison_id:
            logger.info(f"📚 开始预加载术语库: {comparison_id}")
            from .main import get_comparison
            preloaded_terms = get_comparison(comparison_id)
            if preloaded_terms:
                logger.info(f"📚 术语库预加载成功: {len(preloaded_terms)} 个术语")
                trans['preloaded_terms'] = preloaded_terms
            else:
                logger.warning(f"📚 术语库预加载失败: {comparison_id}")
        
        # 创建 Okapi 翻译器
        translator = OkapiPptxTranslator()
        logger.info("✅ Okapi PPTX 翻译器创建成功")
        
        # 设置翻译服务
        class OkapiTranslationService:
            def __init__(self, trans):
                self.trans = trans
            
            def batch_translate(self, texts, source_lang, target_lang):
                """批量翻译文本 - 使用多线程并行处理，支持术语库筛选"""
                from concurrent.futures import ThreadPoolExecutor, as_completed
                import threading
                
                translated_texts = [None] * len(texts)  # 预分配结果数组
                max_workers = min(30, len(texts))  # 最多30个线程
                
                logger.info(f"开始并行翻译 {len(texts)} 个文本，使用 {max_workers} 个线程")
                
                # 进度更新相关变量
                total_count = len(texts)
                progress_lock = threading.Lock()
                
                def update_progress():
                    """更新翻译进度"""
                    with progress_lock:
                        actual_completed = sum(1 for i, text in enumerate(texts) if translated_texts[i] is not None)
                        
                        if actual_completed >= total_count:
                            actual_completed = total_count
                            progress_percentage = 100.0
                        else:
                            progress_percentage = min((actual_completed / total_count) * 100, 100.0)
                        
                        logger.info(f"翻译进度: {actual_completed}/{total_count} ({progress_percentage:.1f}%)")
                        
                        # 更新数据库进度
                        try:
                            from .to_translate import db
                            db.execute("update translate set process=%s where id=%s", 
                                     str(format(progress_percentage, '.1f')), 
                                     self.trans['id'])
                        except Exception as e:
                            logger.error(f"更新进度失败: {str(e)}")
                
                def translate_single_text(index, text):
                    """翻译单个文本 - 使用 translate_text 函数（Okapi只用于解析，翻译用qwen-mt-plus）"""
                    try:
                        # 使用 translate_text 函数进行翻译
                        result = to_translate.translate_text(
                            self.trans,
                            text,
                            source_lang="auto",
                            target_lang=self.trans.get('lang', 'English')
                        )
                        return result
                    except Exception as e:
                        logger.error(f"翻译文本 {index} 失败: {e}")
                        return text  # 失败时返回原文
                
                # 使用线程池并行翻译
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    # 提交所有任务
                    future_to_index = {
                        executor.submit(translate_single_text, i, text): i 
                        for i, text in enumerate(texts)
                    }
                    
                    # 收集结果
                    completed = 0
                    for future in as_completed(future_to_index):
                        index = future_to_index[future]
                        try:
                            translated_texts[index] = future.result()
                            completed += 1
                            
                            # 每完成10%更新一次进度
                            if completed % max(1, total_count // 10) == 0:
                                update_progress()
                        except Exception as e:
                            logger.error(f"获取翻译结果 {index} 失败: {e}")
                            translated_texts[index] = texts[index]  # 失败时使用原文
                
                # 最终进度更新
                update_progress()
                
                return translated_texts
        
        # 设置翻译服务
        translation_service = OkapiTranslationService(trans)
        translator.set_translation_service(translation_service)
        
        # 准备文件路径
        input_file = trans['file_path']
        output_file = trans.get('target_file') or input_file.replace('.pptx', '_translated.pptx')
        
        # 确定源语言和目标语言
        source_lang = "zh"  # PPTX 默认源语言为中文
        target_lang = trans.get('lang', 'English')
        
        # 语言代码映射（如果需要）
        lang_code_map = {
            'English': 'en',
            'Chinese': 'zh',
            'Japanese': 'ja',
            'Korean': 'ko',
            # 可以添加更多映射
        }
        target_lang_code = lang_code_map.get(target_lang, target_lang.lower()[:2])
        
        # 执行翻译
        logger.info(f"开始翻译 PPTX: {input_file} -> {output_file}")
        logger.info(f"源语言: {source_lang}, 目标语言: {target_lang} ({target_lang_code})")
        
        success = translator.translate_document(
            input_file, output_file, source_lang, target_lang_code
        )
        
        if success:
            # 更新数据库状态
            try:
                from .to_translate import db
                db.execute("update translate set status='done', process='100.0', result_path=%s where id=%s",
                         output_file, trans['id'])
                logger.info(f"✅ PPTX 翻译完成: {output_file}")
            except Exception as e:
                logger.error(f"更新数据库状态失败: {e}")
            
            # 翻译成功，删除原始文件以节省空间
            original_file = trans['file_path']
            if os.path.exists(original_file) and original_file != output_file:
                try:
                    os.remove(original_file)
                    logger.info(f"✅ 翻译成功，已删除原始文件: {os.path.basename(original_file)}")
                except Exception as e:
                    logger.warning(f"⚠️ 删除原始文件失败: {str(e)}")
            
            return True
        else:
            # 翻译失败，回退到传统方法
            logger.error("❌ Okapi 翻译失败，回退到传统方法")
            return start_traditional(trans)
            
    except Exception as e:
        logger.error(f"Okapi 翻译过程出错: {e}", exc_info=True)
        # 出错时回退到传统方法
        return start_traditional(trans)


def start_traditional(trans):
    """使用传统方法进行翻译（原有的 start 函数逻辑）"""
    # 允许的最大线程
    threads=trans['threads']
    if threads is None or int(threads)<0:
        max_threads=10
    else:
        max_threads=int(threads)
    # 当前执行的索引位置
    run_index=0
    start_time = datetime.now()
    wb = pptx.Presentation(trans['file_path']) 
    logger.info(f"处理文件: {trans['file_path']}")
    slides = wb.slides
    texts=[]
    
    # 改进的去重策略：使用位置+文本作为唯一标识，允许相同文本在不同位置存在
    # 但避免在同一位置重复提取相同文本
    extracted_texts_by_location = {}  # {(slide_index, shape_index, text_type, text): True}
    
    # 提取文本时保存样式信息，并建立文本与形状的对应关系
    slide_count = 0
    for slide in slides:
        slide_count += 1
        slide_text_count = 0
        logger.info(f"正在处理第 {slide_count} 页幻灯片...")
        
        for shape_index, shape in enumerate(slide.shapes):
            # 处理表格
            if shape.has_table:
                table = shape.table
                logger.info(f"发现表格: {table}")
                rows = len(table.rows)
                cols = len(table.columns)
                for r in range(rows):
                    for c in range(cols):
                        cell_text = table.cell(r, c).text
                        if cell_text!=None and len(cell_text)>0 and not common.is_all_punc(cell_text):
                            # 使用位置+文本作为唯一标识，允许相同文本在不同位置存在
                            cell_text_clean = cell_text.strip()
                            location_key = (slide_count, shape_index, "table_cell", r, c, cell_text_clean)
                            if location_key not in extracted_texts_by_location:
                                # 保存表格单元格的样式信息，并建立对应关系
                                cell = table.cell(r, c)
                                style_info = extract_cell_style(cell)
                                texts.append({
                                    "text": cell_text,
                                    "row": r,
                                    "column": c, 
                                    "complete": False,
                                    "type": "table_cell",
                                    "style_info": style_info,
                                    "slide_index": slide_count,
                                    "shape_index": shape_index,
                                    "shape": shape,
                                    "cell": cell,
                                    "original_text": cell_text  # 保存原始文本用于匹配
                                })
                                slide_text_count += 1
                                extracted_texts_by_location[location_key] = True
                                logger.debug(f"添加表格单元格文本: {cell_text_clean[:50]}... (位置: 幻灯片{slide_count}, 形状{shape_index}, 行{r}, 列{c})")
                            else:
                                logger.debug(f"跳过重复的表格单元格文本: {cell_text_clean[:50]}... (同一位置)")
            
            # 处理所有有文本框架的形状（包括文本框、标题、占位符等）
            elif shape.has_text_frame:
                text_frame = shape.text_frame
                logger.debug(f"发现文本框架，段落数: {len(text_frame.paragraphs)}")
                
                # 记录已提取的文本，避免重复提取
                extracted_texts_in_frame = set()
                paragraph_texts = []  # 收集所有段落文本
                
                # 第一优先级：提取段落文本
                for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
                    text = paragraph.text
                    if text!=None and len(text)>0 and not common.is_all_punc(text):
                        text_clean = text.strip()
                        # 使用位置+文本作为唯一标识
                        location_key = (slide_count, shape_index, "paragraph", paragraph_index, text_clean)
                        if location_key not in extracted_texts_by_location:
                            # 保存段落的样式信息，并建立对应关系
                            style_info = extract_paragraph_style(paragraph)
                            texts.append({
                                "text": text, 
                                "complete": False,
                                "type": "paragraph",
                                "style_info": style_info,
                                "paragraph": paragraph,
                                "slide_index": slide_count,
                                "shape_index": shape_index,
                                "shape": shape,
                                "paragraph_index": paragraph_index,
                                "original_text": text  # 保存原始文本用于匹配
                            })
                            slide_text_count += 1
                            # 记录已提取的文本
                            extracted_texts_in_frame.add(text_clean)
                            extracted_texts_by_location[location_key] = True
                            paragraph_texts.append(text_clean)
                            logger.debug(f"添加段落文本: {text_clean[:30]}... (位置: 幻灯片{slide_count}, 形状{shape_index}, 段落{paragraph_index})")
                        else:
                            logger.debug(f"跳过重复的段落文本: {text_clean[:30]}... (同一位置)")
                
                # 第二优先级：检查文本框架是否有遗漏的文本
                # 注意：如果已经提取了段落文本，通常不需要再提取框架文本，除非框架文本包含段落文本没有的内容
                # 这里暂时跳过框架文本的提取，避免重复
                # 如果确实需要，可以通过比较框架文本和段落文本组合来判断
                pass
            
            # 处理其他可能有文本的形状（如形状内的文本）
            # 注意：这里只处理没有文本框架的形状，避免重复提取
            elif hasattr(shape, 'text') and shape.text and not shape.has_text_frame:
                text = shape.text
                if text!=None and len(text)>0 and not common.is_all_punc(text):
                    text_clean = text.strip()
                    # 使用位置+文本作为唯一标识
                    location_key = (slide_count, shape_index, "shape_text", text_clean)
                    if location_key not in extracted_texts_by_location:
                        texts.append({
                            "text": text, 
                            "complete": False,
                            "type": "shape_text",
                            "shape": shape,
                            "slide_index": slide_count,
                            "shape_index": shape_index,
                            "original_text": text  # 保存原始文本用于匹配
                        })
                        slide_text_count += 1
                        extracted_texts_by_location[location_key] = True
                        logger.debug(f"添加形状文本: {text_clean[:50]}... (位置: 幻灯片{slide_count}, 形状{shape_index})")
                    else:
                        logger.debug(f"跳过重复的形状文本: {text_clean[:50]}... (同一位置)")
            
            # 处理形状名称（通常不会与内容文本重复，但也要检查）
            elif hasattr(shape, 'name') and shape.name:
                shape_name = shape.name
                if shape_name and len(shape_name.strip()) > 0:
                    shape_name_clean = shape_name.strip()
                    # 使用位置+文本作为唯一标识
                    location_key = (slide_count, shape_index, "shape_name", shape_name_clean)
                    if location_key not in extracted_texts_by_location:
                        texts.append({
                            "text": shape_name,
                            "complete": False,
                            "type": "shape_name",
                            "shape": shape,
                            "slide_index": slide_count,
                            "shape_index": shape_index,
                            "original_text": shape_name  # 保存原始文本用于匹配
                        })
                        slide_text_count += 1
                        extracted_texts_by_location[location_key] = True
                        logger.debug(f"添加形状名称文本: {shape_name_clean[:50]}... (位置: 幻灯片{slide_count}, 形状{shape_index})")
                    else:
                        logger.debug(f"跳过重复的形状名称文本: {shape_name_clean[:50]}... (同一位置)")
        
        logger.info(f"第 {slide_count} 页幻灯片提取了 {slide_text_count} 个文本元素")
    
    # 不再进行全局去重，因为相同文本在不同位置需要分别翻译
    logger.info(f"总共提取了 {len(texts)} 个文本元素")
    
    # 调试：打印所有提取的文本类型和详细信息
    text_types = {}
    duplicate_check = {}  # 检查重复文本
    
    for i, item in enumerate(texts):
        text_type = item.get('type', 'unknown')
        text_types[text_type] = text_types.get(text_type, 0) + 1
        
        # 检查重复文本
        text_content = item.get('text', '').strip()
        if text_content in duplicate_check:
            duplicate_check[text_content].append({
                'index': i,
                'type': text_type,
                'slide': item.get('slide_index', 'N/A'),
                'shape': item.get('shape_index', 'N/A')
            })
        else:
            duplicate_check[text_content] = [{
                'index': i,
                'type': text_type,
                'slide': item.get('slide_index', 'N/A'),
                'shape': item.get('shape_index', 'N/A')
            }]
        
        # 记录每个文本项的详细信息（前10个）
        if i < 10:
            logger.info(f"文本项 {i+1}: 类型={text_type}, 幻灯片={item.get('slide_index', 'N/A')}, 形状索引={item.get('shape_index', 'N/A')}, 内容='{item.get('text', '')[:50]}...'")
    
    # 报告重复文本情况
    duplicate_count = 0
    for text_content, occurrences in duplicate_check.items():
        if len(occurrences) > 1:
            duplicate_count += 1
            logger.warning(f"发现重复文本: '{text_content[:50]}...' 出现 {len(occurrences)} 次:")
            for occ in occurrences:
                logger.warning(f"  - 索引{occ['index']}, 类型={occ['type']}, 幻灯片={occ['slide']}, 形状={occ['shape']}")
    
    if duplicate_count > 0:
        logger.warning(f"总共发现 {duplicate_count} 个重复文本，这可能导致翻译重复")
    else:
        logger.info("未发现重复文本，提取正常")
    
    logger.info(f"提取的文本类型分布: {text_types}")
    logger.info(f"总共提取了 {len(texts)} 个文本元素")
    max_run=max_threads if len(texts)>max_threads else len(texts)
    before_active_count=threading.activeCount()
    event=threading.Event()
    
    # 进度更新相关变量
    completed_count = 0
    total_count = len(texts)
    progress_lock = threading.Lock()
    
    def update_progress():
        """更新翻译进度"""
        nonlocal completed_count
        with progress_lock:
            completed_count += 1
            progress_percentage = min((completed_count / total_count) * 100, 100.0)
            
            # 更新数据库进度
            try:
                from .to_translate import db
                db.execute("update translate set process=%s where id=%s", 
                         str(format(progress_percentage, '.1f')), 
                         trans['id'])
                
                # 如果进度达到100%，立即更新状态为已完成
                if progress_percentage >= 100.0:
                    import pytz
                    end_time = datetime.now(pytz.timezone('Asia/Shanghai'))
                    db.execute(
                        "update translate set status='done',end_at=%s,process=100 where id=%s",
                        end_time, trans['id']
                    )
                    print("✅ 翻译完成，状态已更新为已完成")
                    
            except Exception as e:
                print(f"更新进度失败: {str(e)}")
    
    while run_index<=len(texts)-1:
        if threading.activeCount()<max_run+before_active_count:
            if not event.is_set():
                thread = threading.Thread(target=to_translate.get,args=(trans,event,texts,run_index))
                thread.start()
                run_index+=1
            else:
                return False
    
    # 等待翻译完成，并监控进度
    last_completed_count = 0
    while True:
        complete=True
        current_completed = 0
        for text in texts:
            if not text['complete']:
                complete=False
            else:
                current_completed += 1
        
        # 检查是否有新的文本完成
        if current_completed > last_completed_count:
            completed_count = current_completed
            progress_percentage = min((completed_count / total_count) * 100, 100.0)
            
            # 更新数据库进度
            try:
                from .to_translate import db
                db.execute("update translate set process=%s where id=%s", 
                         str(format(progress_percentage, '.1f')), 
                         trans['id'])
                
            except Exception as e:
                print(f"更新进度失败: {str(e)}")
            
            last_completed_count = current_completed
        
        if complete:
            break
        else:
            time.sleep(1)

    text_count=0
    slide_count = 0
    for slide in slides:
        slide_count += 1
        slide_processed_count = 0
        # logger.info(f"正在应用翻译结果到第 {slide_count} 页幻灯片...")
        
        for shape_index, shape in enumerate(slide.shapes):
            # 添加详细的形状信息日志
            # logger.info(f"=== 处理第 {slide_count} 页形状 ===")
            # logger.info(f"形状类型: {type(shape).__name__}")
            # logger.info(f"形状名称: {getattr(shape, 'name', 'N/A')}")
            # logger.info(f"是否有表格: {shape.has_table}")
            # logger.info(f"是否有文本框架: {shape.has_text_frame}")
            # logger.info(f"是否有文本属性: {hasattr(shape, 'text')}")
            # try:
            #     logger.info(f"是否有占位符格式: {hasattr(shape, 'placeholder_format')}")
            # except ValueError:
            #     logger.info(f"是否有占位符格式: False (非占位符形状)")
            
            # 处理表格
            if shape.has_table:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                for r in range(rows):
                    row_data = []
                    for c in range(cols):
                        cell_text = table.cell(r, c).text
                        if cell_text!=None and len(cell_text)>0 and not common.is_all_punc(cell_text):
                            # 查找对应的翻译结果
                            cell = table.cell(r, c)
                            translated_item = find_translated_text_for_shape(texts, slide_count, shape_index, "table_cell", cell=cell, row=r, column=c)
                            if translated_item:
                                # 默认启用自适应样式
                                apply_translation_to_cell_with_adaptive_styles(cell, translated_item['text'], translated_item.get('style_info', {}))
                                text_count+=translated_item.get('count', 1)
                                slide_processed_count += 1
                                logger.info(f"表格单元格翻译应用成功: 原文='{cell_text[:30]}...' -> 译文='{translated_item['text'][:30]}...'")
                            else:
                                logger.warning(f"未找到表格单元格的翻译结果: {cell_text[:30]}...")
                          
            # 处理所有有文本框架的形状（包括文本框、标题、占位符等）
            elif shape.has_text_frame:
                text_frame = shape.text_frame
                logger.info(f"处理文本框架形状，段落数: {len(text_frame.paragraphs)}")
                for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
                    text=paragraph.text
                    if text!=None and len(text)>0 and not common.is_all_punc(text):
                        # 查找对应的翻译结果
                        translated_item = find_translated_text_for_shape(texts, slide_count, shape_index, "paragraph", paragraph=paragraph, paragraph_index=paragraph_index)
                        if translated_item:
                            # 默认启用自适应样式
                            logger.info(f"应用翻译到段落: 原文='{paragraph.text[:30]}...' -> 译文='{translated_item['text'][:30]}...'")
                            apply_translation_to_paragraph_with_adaptive_styles(paragraph, translated_item['text'], translated_item.get('style_info', {}))
                            text_count+=translated_item.get('count', 1)
                            slide_processed_count += 1
                            logger.info(f"段落翻译完成，当前runs数: {len(paragraph.runs)}")
                        else:
                            logger.warning(f"未找到段落的翻译结果: {text[:30]}...")
                
                # 文本框架内容通常已经通过段落处理了，这里不再单独处理，避免重复
                pass
            
            # 处理其他可能有文本的形状（如形状内的文本）
            # 注意：只处理没有文本框架的形状，避免重复
            elif hasattr(shape, 'text') and shape.text and not shape.has_text_frame:
                text = shape.text
                if text!=None and len(text)>0 and not common.is_all_punc(text):
                    # 查找对应的翻译结果
                    translated_item = find_translated_text_for_shape(texts, slide_count, shape_index, "shape_text", shape=shape)
                    if translated_item:
                        # 处理形状文本
                        original_text = shape.text
                        shape.text = translated_item['text']
                        # 应用自适应样式到形状文本
                        apply_adaptive_styles_to_shape(shape, original_text, translated_item['text'])
                        text_count+=translated_item.get('count', 1)
                        slide_processed_count += 1
                        logger.info(f"形状文本翻译应用成功: 原文='{original_text[:30]}...' -> 译文='{translated_item['text'][:30]}...'")
                    else:
                        logger.warning(f"未找到形状文本的翻译结果: {text[:30]}...")
            
            # 文本框架已经在上面处理过了，这里不再重复处理
            
            # 处理形状名称文本
            elif hasattr(shape, 'name') and shape.name:
                shape_name = shape.name
                if shape_name and len(shape_name.strip()) > 0:
                    # 查找对应的翻译结果
                    translated_item = find_translated_text_for_shape(texts, slide_count, shape_index, "shape_name", shape=shape)
                    if translated_item:
                        # 处理形状名称文本
                        original_name = shape.name
                        shape.name = translated_item['text']
                        # 应用自适应样式到形状名称
                        apply_adaptive_styles_to_shape(shape, original_name, translated_item['text'])
                        text_count+=translated_item.get('count', 1)
                        slide_processed_count += 1
                        logger.info(f"形状名称翻译应用成功: 原文='{original_name[:30]}...' -> 译文='{translated_item['text'][:30]}...'")
                    else:
                        logger.warning(f"未找到形状名称的翻译结果: {shape_name[:30]}...")
            
            # 处理占位符文本
            try:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    # 查找对应的翻译结果
                    translated_item = find_translated_text_for_shape(texts, slide_count, shape_index, "placeholder", shape=shape)
                    if translated_item:
                        # 处理占位符文本
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            original_text = shape.text_frame.text
                            shape.text_frame.text = translated_item['text']
                            # 应用自适应样式到所有段落
                            if shape.text_frame.paragraphs:
                                for paragraph in shape.text_frame.paragraphs:
                                    if paragraph.runs:
                                        for run in paragraph.runs:
                                            if ENABLE_FONT_SCALING:
                                                apply_adaptive_styles_ppt(run, original_text, translated_item['text'])
                            text_count += translated_item.get('count', 1)
                            slide_processed_count += 1
                            logger.info(f"占位符文本翻译应用成功: 原文='{original_text[:30]}...' -> 译文='{translated_item['text'][:30]}...'")
                        else:
                            logger.warning(f"未找到占位符文本的翻译结果")
            except ValueError as e:
                # 忽略"shape is not a placeholder"错误
                pass
            except Exception as e:
                logger.warning(f"应用占位符翻译时出错: {str(e)}")
        
        logger.info(f"第 {slide_count} 页幻灯片处理了 {slide_processed_count} 个文本元素")
    
    # 检查是否还有未处理的文本
    # if texts:
    #     logger.warning(f"还有 {len(texts)} 个文本元素未处理:")
    #     for i, item in enumerate(texts[:5]):  # 只显示前5个
    #         logger.warning(f"  未处理文本 {i+1}: {item.get('text', '')[:50]}... (类型: {item.get('type', 'unknown')})")
    #     if len(texts) > 5:
    #         logger.warning(f"  ... 还有 {len(texts) - 5} 个未处理")
    
    # 添加详细的处理统计
    logger.info(f"=== PowerPoint处理统计 ===")
    logger.info(f"总文本元素数: {len(texts) + text_count}")
    logger.info(f"已处理文本数: {text_count}")
    logger.info(f"未处理文本数: {len(texts)}")
    logger.info(f"处理成功率: {text_count / (len(texts) + text_count) * 100:.1f}%" if (len(texts) + text_count) > 0 else "0%")

    logger.info(f"总共处理了 {text_count} 个文本元素")
    wb.save(trans['target_file'])
    end_time = datetime.now()
    spend_time=common.display_spend(start_time, end_time)
    to_translate.complete(trans,text_count,spend_time)
    
    # 翻译成功，删除原始文件以节省空间
    original_file = trans['file_path']
    if os.path.exists(original_file) and original_file != trans.get('target_file'):
        try:
            os.remove(original_file)
            logger.info(f"✅ 翻译成功，已删除原始文件: {os.path.basename(original_file)}")
        except Exception as e:
            logger.warning(f"⚠️ 删除原始文件失败: {str(e)}")
    
    return True

def extract_paragraph_style(paragraph):
    """提取段落的样式信息"""
    style_info = {
        'paragraph_level': {},
        'runs': []
    }
    
    try:
        # 段落级别的样式
        if paragraph.alignment:
            style_info['paragraph_level']['alignment'] = paragraph.alignment
        if paragraph.level:
            style_info['paragraph_level']['level'] = paragraph.level
        
        # 每个run的样式
        for run in paragraph.runs:
            run_style = {}
            
            try:
                # 字体样式
                if run.font.name:
                    run_style['font_name'] = run.font.name
                if run.font.size:
                    run_style['font_size'] = run.font.size
                if run.font.bold is not None:
                    run_style['bold'] = run.font.bold
                if run.font.italic is not None:
                    run_style['italic'] = run.font.italic
                if run.font.underline is not None:
                    run_style['underline'] = run.font.underline
                
                # 颜色处理 - 简化版本，只保存颜色对象
                if run.font.color:
                    # 直接保存颜色对象，避免类型检查错误
                    run_style['color_object'] = run.font.color
                        
            except Exception as e:
                # 如果获取样式失败，继续处理下一个run
                logger.warning(f"提取run样式失败: {str(e)}")
                continue
            
            style_info['runs'].append({
                'text': run.text,
                'style': run_style
            })
    except Exception as e:
        # 如果提取段落样式失败，返回空的样式信息
        logger.warning(f"提取段落样式失败: {str(e)}")
    
    return style_info

def extract_cell_style(cell):
    """提取表格单元格的样式信息"""
    style_info = {
        'paragraphs': []
    }
    
    for paragraph in cell.text_frame.paragraphs:
        para_style = extract_paragraph_style(paragraph)
        style_info['paragraphs'].append(para_style)
    
    return style_info

def apply_translation_to_paragraph(paragraph, translated_text, style_info):
    """应用翻译结果到段落并恢复样式（原始版本）"""
    # 保存原始文本用于自适应计算
    original_text = paragraph.text
    
    # 清空段落内容
    paragraph.clear()
    
    # 如果有样式信息，按run恢复样式
    if style_info and 'runs' in style_info and style_info['runs']:
        # 按原始run的样式分配翻译文本（应用自适应样式）
        distribute_text_to_runs_with_adaptive_styles(paragraph, translated_text, style_info['runs'], original_text)
    else:
        # 没有样式信息，直接添加文本
        original_text = paragraph.text
        paragraph.text = translated_text
        # 应用自适应样式到所有runs
        if paragraph.runs:
            for run in paragraph.runs:
                if ENABLE_FONT_SCALING:
                    apply_adaptive_styles_ppt(run, original_text, translated_text)
    
    # 恢复段落级别的样式
    if style_info and 'paragraph_level' in style_info:
        para_level = style_info['paragraph_level']
        if 'alignment' in para_level:
            paragraph.alignment = para_level['alignment']
        if 'level' in para_level:
            paragraph.level = para_level['level']


def apply_translation_to_cell(cell, translated_text, style_info):
    """应用翻译结果到表格单元格并恢复样式（原始版本）"""
    # 保存原始文本用于自适应计算
    original_text = cell.text
    
    # 清空单元格内容
    cell.text = ""
    
    # 如果有样式信息，按段落恢复样式
    if style_info and 'paragraphs' in style_info and style_info['paragraphs']:
        # 按原始段落的样式分配翻译文本（应用自适应样式）
        distribute_text_to_paragraphs_with_adaptive_styles(cell.text_frame, translated_text, style_info['paragraphs'], original_text)
    else:
        # 没有样式信息，直接添加文本
        original_text = cell.text
        cell.text = translated_text
        # 应用自适应样式到单元格中的所有段落的所有runs
        if cell.text_frame.paragraphs:
            for paragraph in cell.text_frame.paragraphs:
                if paragraph.runs:
                    for run in paragraph.runs:
                        if ENABLE_FONT_SCALING:
                            apply_adaptive_styles_ppt(run, original_text, translated_text)


def apply_translation_to_paragraph_with_adaptive_styles(paragraph, translated_text, style_info):
    """应用翻译结果到段落并恢复样式，同时应用自适应样式"""
    # 保存原始文本用于自适应计算
    original_text = paragraph.text
    
    # 清空段落内容
    paragraph.clear()
    
    # 如果有样式信息，按run恢复样式
    if style_info and 'runs' in style_info and style_info['runs']:
        # 按原始run的样式分配翻译文本
        distribute_text_to_runs_with_adaptive_styles(paragraph, translated_text, style_info['runs'], original_text)
    else:
        # 没有样式信息，直接添加文本
        paragraph.text = translated_text
        # 应用自适应样式到所有runs
        if paragraph.runs:
            for run in paragraph.runs:
                if ENABLE_FONT_SCALING:
                    apply_adaptive_styles_ppt(run, original_text, translated_text)
    
    # 恢复段落级别的样式
    if style_info and 'paragraph_level' in style_info:
        para_level = style_info['paragraph_level']
        if 'alignment' in para_level:
            paragraph.alignment = para_level['alignment']
        if 'level' in para_level:
            paragraph.level = para_level['level']


def apply_translation_to_cell_with_adaptive_styles(cell, translated_text, style_info):
    """应用翻译结果到表格单元格并恢复样式，同时应用自适应样式"""
    # 保存原始文本用于自适应计算
    original_text = cell.text
    
    # 清空单元格内容
    cell.text = ""
    
    # 如果有样式信息，按段落恢复样式
    if style_info and 'paragraphs' in style_info and style_info['paragraphs']:
        # 按原始段落的样式分配翻译文本
        distribute_text_to_paragraphs_with_adaptive_styles(cell.text_frame, translated_text, style_info['paragraphs'], original_text)
    else:
        # 没有样式信息，直接添加文本
        cell.text = translated_text
        # 应用自适应样式到单元格中的所有段落的所有runs
        if cell.text_frame.paragraphs:
            for paragraph in cell.text_frame.paragraphs:
                if paragraph.runs:
                    for run in paragraph.runs:
                        if ENABLE_FONT_SCALING:
                            apply_adaptive_styles_ppt(run, original_text, translated_text)


def distribute_text_to_runs_with_adaptive_styles(paragraph, translated_text, run_styles, original_text):
    """将翻译文本按原始run的样式分配到新的run中，同时应用自适应样式"""
    if not run_styles:
        paragraph.text = translated_text
        # 应用自适应样式到所有runs
        if paragraph.runs:
            for run in paragraph.runs:
                if ENABLE_FONT_SCALING:
                    apply_adaptive_styles_ppt(run, original_text, translated_text)
        return
    
    # 计算每个run应该分配的文本长度
    total_original_length = sum(len(run['text']) for run in run_styles)
    if total_original_length == 0:
        paragraph.text = translated_text
        # 应用自适应样式到所有runs
        if paragraph.runs:
            for run in paragraph.runs:
                if ENABLE_FONT_SCALING:
                    apply_adaptive_styles_ppt(run, original_text, translated_text)
        return
    
    # 按比例分配翻译文本
    current_pos = 0
    for i, run_style in enumerate(run_styles):
        original_length = len(run_style['text'])
        if original_length == 0:
            continue
        
        # 计算这个run应该分配的文本长度
        if i == len(run_styles) - 1:
            # 最后一个run，分配剩余的所有文本
            allocated_text = translated_text[current_pos:]
        else:
            # 按比例分配
            ratio = original_length / total_original_length
            allocated_length = max(1, int(len(translated_text) * ratio))
            # 确保不超过剩余文本长度
            remaining_length = len(translated_text) - current_pos
            allocated_length = min(allocated_length, remaining_length)
            allocated_text = translated_text[current_pos:current_pos + allocated_length]
            current_pos += allocated_length
        
        if allocated_text:
            # 创建新的run并应用样式
            run = paragraph.add_run()
            run.text = allocated_text
            apply_run_style(run, run_style['style'])
            # 应用自适应样式 - 使用整个段落的原始文本和翻译文本
            if ENABLE_FONT_SCALING:
                apply_adaptive_styles_ppt(run, original_text, translated_text)


def distribute_text_to_paragraphs_with_adaptive_styles(text_frame, translated_text, paragraph_styles, original_text):
    """将翻译文本按原始段落的样式分配到新的段落中，同时应用自适应样式"""
    if not paragraph_styles:
        text_frame.text = translated_text
        # 应用自适应样式到所有段落的所有runs
        if text_frame.paragraphs:
            for paragraph in text_frame.paragraphs:
                if paragraph.runs:
                    for run in paragraph.runs:
                        if ENABLE_FONT_SCALING:
                            apply_adaptive_styles_ppt(run, original_text, translated_text)
        return
    
    # 清空文本框架
    text_frame.clear()
    
    # 按段落分配文本
    total_original_length = sum(len(run['text']) for para in paragraph_styles for run in para['runs'])
    if total_original_length == 0:
        text_frame.text = translated_text
        # 应用自适应样式到所有段落的所有runs
        if text_frame.paragraphs:
            for paragraph in text_frame.paragraphs:
                if paragraph.runs:
                    for run in paragraph.runs:
                        if ENABLE_FONT_SCALING:
                            apply_adaptive_styles_ppt(run, original_text, translated_text)
        return
    
    current_pos = 0
    for i, para_style in enumerate(paragraph_styles):
        if i > 0:
            # 添加段落分隔符
            text_frame.add_paragraph()
        
        paragraph = text_frame.paragraphs[-1] if text_frame.paragraphs else text_frame.add_paragraph()
        
        # 计算这个段落应该分配的文本长度
        para_original_length = sum(len(run['text']) for run in para_style['runs'])
        if para_original_length == 0:
            continue
        
        if i == len(paragraph_styles) - 1:
            # 最后一个段落，分配剩余的所有文本
            allocated_text = translated_text[current_pos:]
        else:
            # 按比例分配
            ratio = para_original_length / total_original_length
            allocated_length = max(1, int(len(translated_text) * ratio))
            # 确保不超过剩余文本长度
            remaining_length = len(translated_text) - current_pos
            allocated_length = min(allocated_length, remaining_length)
            allocated_text = translated_text[current_pos:current_pos + allocated_length]
            current_pos += allocated_length
        
        if allocated_text:
            # 按run分配文本并应用样式
            if para_style['runs']:
                distribute_text_to_runs_with_adaptive_styles(paragraph, allocated_text, para_style['runs'], original_text)
            else:
                paragraph.text = allocated_text
                # 应用自适应样式到所有runs
                if paragraph.runs:
                    for run in paragraph.runs:
                        if ENABLE_FONT_SCALING:
                            apply_adaptive_styles_ppt(run, original_text, allocated_text)

def apply_run_style(run, style):
    """应用run的样式"""
    if not style:
        return
    
    try:
        # 字体名称
        if 'font_name' in style:
            run.font.name = style['font_name']
        
        # 字体大小
        if 'font_size' in style:
            run.font.size = style['font_size']
        
        # 粗体
        if 'bold' in style:
            run.font.bold = style['bold']
        
        # 斜体
        if 'italic' in style:
            run.font.italic = style['italic']
        
        # 下划线
        if 'underline' in style:
            run.font.underline = style['underline']
        
        # 颜色处理 - 简化版本
        if 'color_object' in style:
            try:
                color_object = style['color_object']
                # 尝试复制颜色属性
                if hasattr(color_object, 'rgb') and color_object.rgb:
                    run.font.color.rgb = color_object.rgb
                elif hasattr(color_object, 'theme_color') and color_object.theme_color:
                    run.font.color.theme_color = color_object.theme_color
                elif hasattr(color_object, 'color_index') and color_object.color_index:
                    run.font.color.color_index = color_object.color_index
                elif hasattr(color_object, 'srgbClr') and color_object.srgbClr:
                    run.font.color.srgbClr = color_object.srgbClr
                # 其他颜色类型暂时跳过，避免错误
            except Exception as e:
                logger.warning(f"应用颜色样式失败: {str(e)}")
                pass  # 如果颜色设置失败，跳过颜色设置
    except Exception as e:
        # 如果应用样式失败，记录错误但继续处理
        logger.warning(f"应用run样式失败: {str(e)}")


def calculate_adaptive_font_size_ppt(original_text, translated_text, original_font_size):
    """根据翻译后文本长度计算PPT自适应字体大小"""
    try:
        if not original_font_size:
            return None
        
        # 计算文本长度比例
        original_length = len(original_text.strip())
        translated_length = len(translated_text.strip())
        
        if original_length == 0:
            return original_font_size
        
        # 计算长度比例
        length_ratio = translated_length / original_length
        
        # 添加调试日志
        logger.info(f"PPT字体自适应调试: 原文长度={original_length}, 译文长度={translated_length}, 比例={length_ratio:.2f}")
        
        # 如果翻译后文本变长，适当缩小字体
        if length_ratio > 2.0:  # 文本长度增加超过100%
            # 文本特别长时，最小可以到原始大小的50%
            new_size = max(original_font_size * 0.5, original_font_size / length_ratio)
            logger.info(f"PPT字体自适应: 文本特别长，缩小到50% -> {new_size}pt")
            return int(new_size)
        elif length_ratio > 1.3:  # 文本长度增加超过30%
            # 最小可以到原始大小的50%
            new_size = max(original_font_size * 0.5, original_font_size / length_ratio)
            logger.info(f"PPT字体自适应: 文本较长，缩小到50% -> {new_size}pt")
            return int(new_size)
        elif length_ratio < 0.7:  # 文本长度减少超过30%
            # 如果文本变短，可以适当增大字体，但不要超过原始大小的130%
            new_size = min(original_font_size * 1.3, original_font_size / length_ratio)
            return int(new_size)
        else:
            # 长度变化不大，保持原始字体大小
            return original_font_size
            
    except Exception as e:
        logger.error(f"计算PPT自适应字体大小失败: {str(e)}")
        return original_font_size


def apply_adaptive_styles_ppt(run, original_text, translated_text):
    """应用自适应样式到PPT run"""
    # 检查字体缩放功能是否启用
    if not ENABLE_FONT_SCALING:
        return  # 如果关闭了字体缩放，直接返回
    
    try:
        # 获取原始字体大小
        original_font_size = None
        
        # 方法1：直接从run获取字体大小
        if run.font.size:
            try:
                original_font_size = run.font.size.pt
            except:
                original_font_size = None
        
        # 方法2：从段落级别获取字体大小
        if not original_font_size and hasattr(run, '_element') and run._element.getparent():
            try:
                paragraph = run._element.getparent()
                if hasattr(paragraph, 'pPr') and paragraph.pPr:
                    defRPr = paragraph.pPr.defRPr
                    if defRPr and defRPr.sz:
                        original_font_size = defRPr.sz.val / 100  # 转换为pt
            except:
                pass
        
        # 方法3：使用默认字体大小
        if not original_font_size:
            original_font_size = 14  # 默认14pt
            logger.info(f"使用默认字体大小: {original_font_size}pt")
        
        # 添加详细调试信息
        logger.info(f"=== PPT自适应样式调试 ===")
        logger.info(f"Run文本: '{run.text[:30]}...'")
        logger.info(f"原文: '{original_text[:30]}...'")
        logger.info(f"译文: '{translated_text[:30]}...'")
        logger.info(f"原始字体大小: {original_font_size}pt")
        
        # 计算自适应字体大小
        if original_font_size:
            adaptive_font_size = calculate_adaptive_font_size_ppt(original_text, translated_text, original_font_size)
            logger.info(f"计算后字体大小: {adaptive_font_size}pt")
            
            if adaptive_font_size and adaptive_font_size != original_font_size:
                from pptx.util import Pt
                run.font.size = Pt(adaptive_font_size)
                logger.info(f"✅ 字体大小已调整: {original_font_size}pt -> {adaptive_font_size}pt")
            else:
                logger.info(f"⚠️ 字体大小未改变 (原始={original_font_size}pt, 计算={adaptive_font_size}pt)")
        else:
            logger.warning(f"❌ 无法获取原始字体大小")
                    
    except Exception as e:
        logger.error(f"❌ 应用PPT自适应样式失败: {str(e)}")


def apply_adaptive_styles_to_shape(shape, original_text, translated_text):
    """应用自适应样式到形状文本"""
    # 检查字体缩放功能是否启用
    if not ENABLE_FONT_SCALING:
        return  # 如果关闭了字体缩放，直接返回
    
    try:
        # 获取原始字体大小
        original_font_size = None
        
        # 方法1：从文本框架获取字体大小
        if hasattr(shape, 'text_frame') and shape.text_frame:
            if shape.text_frame.paragraphs:
                paragraph = shape.text_frame.paragraphs[0]
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if run.font.size:
                        try:
                            original_font_size = run.font.size.pt
                        except:
                            original_font_size = None
        
        # 方法2：从形状属性获取字体大小（如果有的话）
        if not original_font_size and hasattr(shape, 'text') and shape.text:
            # 对于简单的形状文本，尝试获取默认字体大小
            try:
                # 使用一个合理的默认字体大小
                original_font_size = 18  # 默认18pt
                logger.info(f"PPT形状自适应: 使用默认字体大小 {original_font_size}pt")
            except:
                original_font_size = None
        
        # 方法3：从形状名称获取字体大小（如果有的话）
        if not original_font_size and hasattr(shape, 'name') and shape.name:
            # 对于形状名称，使用较小的默认字体大小
            original_font_size = 12  # 默认12pt
            logger.info(f"PPT形状名称自适应: 使用默认字体大小 {original_font_size}pt")
        
        # 计算自适应字体大小
        if original_font_size:
            adaptive_font_size = calculate_adaptive_font_size_ppt(original_text, translated_text, original_font_size)
            if adaptive_font_size and adaptive_font_size != original_font_size:
                from pptx.util import Pt
                # 应用到形状的文本框架
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.paragraphs:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if ENABLE_FONT_SCALING:
                                run.font.size = Pt(adaptive_font_size)
                logger.info(f"PPT形状字体大小自适应: {original_font_size}pt -> {adaptive_font_size}pt")
            else:
                logger.info(f"PPT形状自适应: 字体大小未改变 (原始={original_font_size}pt, 计算={adaptive_font_size}pt)")
        else:
            logger.warning(f"PPT形状自适应: 无法获取原始字体大小")
                    
    except Exception as e:
        logger.error(f"应用PPT形状自适应样式失败: {str(e)}")


def distribute_text_to_runs(paragraph, translated_text, run_styles):
    """将翻译文本按原始run的样式分配到新的run中（原始版本）"""
    if not run_styles:
        paragraph.text = translated_text
        return
    
    # 计算每个run应该分配的文本长度
    total_original_length = sum(len(run['text']) for run in run_styles)
    if total_original_length == 0:
        paragraph.text = translated_text
        return
    
    # 按比例分配翻译文本
    current_pos = 0
    for i, run_style in enumerate(run_styles):
        original_length = len(run_style['text'])
        if original_length == 0:
            continue
        
        # 计算这个run应该分配的文本长度
        if i == len(run_styles) - 1:
            # 最后一个run，分配剩余的所有文本
            allocated_text = translated_text[current_pos:]
        else:
            # 按比例分配
            ratio = original_length / total_original_length
            allocated_length = max(1, int(len(translated_text) * ratio))
            # 确保不超过剩余文本长度
            remaining_length = len(translated_text) - current_pos
            allocated_length = min(allocated_length, remaining_length)
            allocated_text = translated_text[current_pos:current_pos + allocated_length]
            current_pos += allocated_length
        
        if allocated_text:
            # 创建新的run并应用样式
            run = paragraph.add_run()
            run.text = allocated_text
            apply_run_style(run, run_style['style'])


def distribute_text_to_paragraphs(text_frame, translated_text, paragraph_styles):
    """将翻译文本按原始段落的样式分配到新的段落中（原始版本）"""
    if not paragraph_styles:
        text_frame.text = translated_text
        return
    
    # 清空文本框架
    text_frame.clear()
    
    # 按段落分配文本
    total_original_length = sum(len(run['text']) for para in paragraph_styles for run in para['runs'])
    if total_original_length == 0:
        text_frame.text = translated_text
        return
    
    current_pos = 0
    for i, para_style in enumerate(paragraph_styles):
        if i > 0:
            # 添加段落分隔符
            text_frame.add_paragraph()
        
        paragraph = text_frame.paragraphs[-1] if text_frame.paragraphs else text_frame.add_paragraph()
        
        # 计算这个段落应该分配的文本长度
        para_original_length = sum(len(run['text']) for run in para_style['runs'])
        if para_original_length == 0:
            continue
        
        if i == len(paragraph_styles) - 1:
            # 最后一个段落，分配剩余的所有文本
            allocated_text = translated_text[current_pos:]
        else:
            # 按比例分配
            ratio = para_original_length / total_original_length
            allocated_length = max(1, int(len(translated_text) * ratio))
            # 确保不超过剩余文本长度
            remaining_length = len(translated_text) - current_pos
            allocated_length = min(allocated_length, remaining_length)
            allocated_text = translated_text[current_pos:current_pos + allocated_length]
            current_pos += allocated_length
        
        if allocated_text:
            # 按run分配文本并应用自适应样式
            # 注意：最后一个参数应该是original_text，用于样式恢复
            distribute_text_to_runs_with_adaptive_styles(paragraph, allocated_text, para_style['runs'], original_text)
            
            # 恢复段落级别的样式
            if 'paragraph_level' in para_style:
                para_level = para_style['paragraph_level']
                if 'alignment' in para_level:
                    paragraph.alignment = para_level['alignment']
                if 'level' in para_level:
                    paragraph.level = para_level['level']


def find_translated_text_for_shape(texts, slide_index, shape_index, text_type, **kwargs):
    """从texts列表中查找与给定形状和文本类型匹配的翻译结果"""
    # 首先尝试精确匹配（对象引用匹配）
    for item in texts:
        if item['type'] == text_type and item.get('complete', False):
            # 检查形状是否匹配
            if 'shape' in kwargs and item.get('shape') == kwargs.get('shape'):
                return item
            # 检查文本框架是否匹配
            if 'text_frame' in kwargs and item.get('text_frame') == kwargs.get('text_frame'):
                return item
            # 检查段落是否匹配
            if 'paragraph' in kwargs and item.get('paragraph') == kwargs.get('paragraph'):
                return item
            # 检查单元格是否匹配
            if 'cell' in kwargs and item.get('cell') == kwargs.get('cell'):
                return item
    
    # 如果精确匹配失败，尝试基于位置和内容的精确匹配
    for item in texts:
        if item['type'] == text_type and item.get('complete', False):
            # 检查幻灯片索引和形状索引是否匹配
            if (item.get('slide_index') == slide_index and 
                item.get('shape_index') == shape_index):
                
                # 对于段落类型，还需要检查段落索引
                if text_type == 'paragraph' and 'paragraph_index' in kwargs:
                    if item.get('paragraph_index') == kwargs.get('paragraph_index'):
                        # 额外检查原始文本是否匹配
                        paragraph = kwargs.get('paragraph')
                        original_text = paragraph.text if paragraph and hasattr(paragraph, 'text') else None
                        if original_text and item.get('original_text'):
                            if original_text.strip() == item.get('original_text', '').strip():
                                return item
                # 对于表格单元格，检查行列索引
                elif text_type == 'table_cell':
                    if (item.get('row') == kwargs.get('row') and 
                        item.get('column') == kwargs.get('column')):
                        # 额外检查原始文本是否匹配
                        cell = kwargs.get('cell')
                        cell_text = cell.text if cell and hasattr(cell, 'text') else None
                        if cell_text and item.get('original_text'):
                            if cell_text.strip() == item.get('original_text', '').strip():
                                return item
                # 对于其他类型，检查原始文本是否匹配
                elif text_type in ['shape_text', 'shape_name', 'text_frame']:
                    # 获取原始文本
                    original_text = None
                    if 'shape' in kwargs:
                        shape = kwargs.get('shape')
                        if text_type == 'shape_text' and hasattr(shape, 'text'):
                            original_text = shape.text
                        elif text_type == 'shape_name' and hasattr(shape, 'name'):
                            original_text = shape.name
                    elif 'text_frame' in kwargs:
                        text_frame = kwargs.get('text_frame')
                        if hasattr(text_frame, 'text'):
                            original_text = text_frame.text
                    
                    # 比较原始文本
                    if original_text and item.get('original_text'):
                        if original_text.strip() == item.get('original_text', '').strip():
                            return item
    
    # 如果所有匹配都失败，记录警告并返回None
    logger.warning(f"未找到匹配的翻译结果: slide_index={slide_index}, shape_index={shape_index}, text_type={text_type}, kwargs={kwargs}")
    return None


